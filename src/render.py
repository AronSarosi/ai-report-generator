"""Render a Report into matplotlib charts + a consulting-grade .pptx, then export .pdf.

Visuals follow docs/report_design_spec.md via src/style.py (one source of truth). Charts
are matplotlib PNGs (full Tufte control, stable in PDF) placed into the python-pptx deck.
PDF export uses LibreOffice headless if present; otherwise the .pptx is still produced.
"""

from __future__ import annotations

import os
import shutil
import subprocess
from collections import namedtuple
from pathlib import Path
from typing import Optional

from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.config import get_settings
from src.report import money, signed_money
from src.schemas import ChartSpec, Report
from src.style import (
    ACCENT,
    CONTENT_W_IN,
    FONT_BODY,
    FONT_HEAD,
    GRAY_100,
    GRAY_300,
    GRAY_500,
    GRAY_900,
    INK,
    KICKER_BOX,
    MARGIN_IN,
    PAPER,
    SLIDE_H_IN,
    SLIDE_W_IN,
    SOURCE_BOX,
    SZ_DECK_TITLE,
    SZ_KICKER,
    SZ_SOURCE,
    SZ_SUB,
    SZ_TITLE,
    TITLE_BOX,
    apply_mpl_style,
)

# --------------------------------------------------------------------------- #
# Brand theme (defaults to the house style; an uploaded template overrides it)
# --------------------------------------------------------------------------- #
Theme = namedtuple("Theme", "accent ink font_head font_body")


def _theme(brand: Optional[dict]) -> Theme:
    brand = brand or {}
    return Theme(accent=brand.get("accent") or ACCENT, ink=brand.get("ink") or INK,
                 font_head=brand.get("font_head") or FONT_HEAD,
                 font_body=brand.get("font_body") or FONT_BODY)


# --------------------------------------------------------------------------- #
# Charts
# --------------------------------------------------------------------------- #
def render_chart(spec: ChartSpec, path: Path, accent: str = ACCENT, wide: bool = False) -> Path:
    import matplotlib.pyplot as plt

    apply_mpl_style()
    fig, ax = plt.subplots(figsize=(10.8, 3.2) if wide else (7.8, 4.0))
    name, ys = next(iter(spec.series.items()))
    labels = spec.x

    if spec.kind == "line":
        xs = list(range(len(spec.x)))
        ax.plot(xs, ys, color=accent, linewidth=2, marker="o", markersize=3)
        if ys:
            ax.scatter([xs[-1]], [ys[-1]], color=accent, s=45, zorder=5)
            ax.annotate(money(ys[-1]), (xs[-1], ys[-1]), textcoords="offset points",
                        xytext=(-4, 8), ha="right", color=INK, fontsize=9, fontweight="bold")
        ax.set_xticks(xs)
        ax.set_xticklabels(spec.x, rotation=45, ha="right", fontsize=8)
        ax.yaxis.set_major_formatter(FuncFormatter(lambda v, p: money(v)))
        ax.margins(x=0.02)
    elif spec.kind == "donut":
        # Share of total: top slices + an "Other" remainder; the leader is accented.
        pairs = list(zip(labels, ys))[:6]
        rest = sum(ys[6:])
        if rest > 0:
            pairs.append(("Other", rest))
        dlabels = [p[0] for p in pairs]
        dvals = [p[1] for p in pairs]
        palette = [accent, INK, GRAY_500, GRAY_300, "#9AA8Bd", "#5B6B82", GRAY_100]
        colors = [accent if (spec.highlight and dlabels[i] == spec.highlight) else
                  palette[i % len(palette)] for i in range(len(dlabels))]
        ax.pie(dvals, labels=dlabels, colors=colors, startangle=90, counterclock=False,
               wedgeprops={"width": 0.42, "edgecolor": PAPER, "linewidth": 2},
               textprops={"fontsize": 9, "color": GRAY_900},
               autopct=lambda p: f"{p:.0f}%" if p >= 6 else "")
        ax.set(aspect="equal")
    elif spec.kind == "column":  # vertical bars
        xs = list(range(len(labels)))
        colors = [accent if (spec.highlight and labels[i] == spec.highlight) else GRAY_300
                  for i in range(len(labels))]
        ax.bar(xs, ys, color=colors)
        ax.set_xticks(xs)
        ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
        ax.yaxis.set_major_formatter(FuncFormatter(lambda v, p: money(v)))
        for x, val in zip(xs, ys):
            ax.annotate(money(val), (x, val), textcoords="offset points",
                        xytext=(0, 4), ha="center", fontsize=8, color=GRAY_900)
    else:  # horizontal bar (sorted desc -> largest on top)
        y_pos = list(range(len(labels)))[::-1]
        colors = [accent if (spec.highlight and labels[i] == spec.highlight) else GRAY_300
                  for i in range(len(labels))]
        ax.barh(y_pos, ys, color=colors)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=9)
        ax.xaxis.set_major_formatter(FuncFormatter(lambda v, p: money(v)))
        ax.grid(False)
        ax.xaxis.grid(True, color=GRAY_100, linewidth=0.8)
        for yp, val in zip(y_pos, ys):
            ax.annotate(money(val), (val, yp), textcoords="offset points",
                        xytext=(4, 0), va="center", fontsize=8, color=GRAY_900)

    fig.tight_layout()
    path = Path(path)
    fig.savefig(path)
    plt.close(fig)
    return path


# --------------------------------------------------------------------------- #
# PPTX helpers
# --------------------------------------------------------------------------- #
def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size, font=FONT_BODY, color=GRAY_900, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.color.rgb = _rgb(color)
    return r


def _rule(slide, left, top, width, pt_height, color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Pt(pt_height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _content_slide(prs, kicker, title, source, t):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set(_box(s, *KICKER_BOX).paragraphs[0], kicker.upper(), SZ_KICKER, t.font_body, t.accent, bold=True)
    _set(_box(s, *TITLE_BOX).paragraphs[0], title, SZ_TITLE, t.font_head, t.ink, bold=True)
    # No rule under the title — whitespace separates it (underlines read as AI-generated).
    if source:
        _set(_box(s, *SOURCE_BOX).paragraphs[0], source, SZ_SOURCE, t.font_body, GRAY_500)
    return s


def _full_bg(slide, color: str):
    """A full-bleed background rectangle (a 'cover' panel)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(color)
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def _cover_slide(prs, report: Report, t, dark: bool = True):
    """The cover: title at the very top with nothing above it, an accent rule, then the
    subtitle. Dark = full-bleed brand-ink panel with a white title; light = clean white
    cover with an ink title (varied across decks)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        _full_bg(s, t.ink)
        title_color, sub_color = PAPER, GRAY_300
    else:
        title_color, sub_color = t.ink, GRAY_500
    _set(_box(s, 0.85, 0.85, 11.6, 2.0).paragraphs[0], report.title, SZ_DECK_TITLE,
         t.font_head, title_color, bold=True)
    _rule(s, 0.9, 3.0, 1.7, 6, color=t.accent)
    _set(_box(s, 0.85, 3.25, 11.6, 0.6).paragraphs[0], report.subtitle, SZ_SUB,
         t.font_body, sub_color)


def _exec_slide(prs, report: Report, t):
    s = _content_slide(prs, "EXECUTIVE SUMMARY", report.governing_thought or "Executive summary",
                       "Source: figures verified against the source data.", t)
    tf = _box(s, 0.5, 2.05, CONTENT_W_IN, 4.6)
    for i, km in enumerate(report.key_messages):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # One universal bullet color (the brand accent) regardless of sentiment — colored
        # red/green dots read like a status alert rather than a finished report.
        chip = p.add_run()
        chip.text = "●  "
        chip.font.color.rgb = _rgb(t.accent)
        chip.font.size = Pt(16)
        chip.font.bold = True
        _set(p, km.text, 16, t.font_body, GRAY_900)
        p.space_after = Pt(16)


def _insight_slide(prs, sec, chart_path: Optional[Path], t):
    src = f"Source: {sec.citations[0] if sec.citations else 'database'}, verified against the data."
    s = _content_slide(prs, sec.kicker, sec.action_title, src, t)
    if chart_path:
        s.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.95), width=Inches(7.6))
    if sec.narrative:
        _set(_box(s, 0.5, 6.0, 7.6, 0.95).paragraphs[0], sec.narrative, 11, t.font_body, GRAY_500)
    _callout(s, "KEY TAKEAWAY", sec.so_what or "", sec.bullets[:3], t)


def _callout(slide, label, takeaway, bullets, t, left=8.5, top=1.95, width=4.3):
    """A brand-colored header band over a light card — the insight callout."""
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.5))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = _rgb(t.ink)
    hdr.line.fill.background()
    hdr.shadow.inherit = False
    htf = hdr.text_frame
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    htf.margin_left = Inches(0.22)
    htf.margin_top = Inches(0.02)
    htf.margin_bottom = Inches(0.02)
    _set(htf.paragraphs[0], label, 12, t.font_body, PAPER, bold=True)

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.5),
                                  Inches(width), Inches(3.5))
    body.fill.solid()
    body.fill.fore_color.rgb = _rgb(GRAY_100)
    body.line.color.rgb = _rgb(GRAY_300)
    body.line.width = Pt(0.75)
    body.shadow.inherit = False
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.26)
    _set(tf.paragraphs[0], takeaway, 13, t.font_body, GRAY_900)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].space_after = Pt(16)
    for b in bullets:
        p = tf.add_paragraph()
        _set(p, "•  " + b, 11, t.font_body, GRAY_900)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(7)


def _kpis_from_chart(spec: ChartSpec) -> list[tuple[str, str]]:
    """Up to three grounded headline numbers (value, label) drawn from the chart data.
    Bars -> the top three; a trend line -> latest, peak, and change vs the start."""
    if not spec or not spec.series:
        return []
    _, ys = next(iter(spec.series.items()))
    xs = spec.x or []
    if not ys:
        return []
    if spec.kind == "line":
        out = [(money(ys[-1]), f"Latest ({xs[-1]})" if xs else "Latest")]
        peak_i = max(range(len(ys)), key=lambda i: ys[i])
        out.append((money(ys[peak_i]), f"Peak ({xs[peak_i]})" if xs else "Peak"))
        out.append((signed_money(ys[-1] - ys[0]), "Change vs start"))
        return out
    pairs = sorted(zip(xs, ys), key=lambda p: p[1], reverse=True)[:3]
    return [(money(v), str(lbl)) for lbl, v in pairs]


def _kpi_band(slide, kpis: list[tuple[str, str]], descriptions: list[str], t, top: float = 5.2):
    """A horizontal row of KPI cards: a big number, then a one-line insight beneath it
    (not just a label — the label is already on the chart above)."""
    n = len(kpis) or 1
    gap = 0.5
    card_w = (CONTENT_W_IN - gap * (n - 1)) / n
    for i, (val, label) in enumerate(kpis):
        x = MARGIN_IN + i * (card_w + gap)
        desc = descriptions[i] if i < len(descriptions) else label
        if len(desc) > 120:
            desc = desc[:117] + "..."
        _rule(slide, x, top, card_w, 4, color=t.accent)
        _set(_box(slide, x, top + 0.1, card_w, 0.7).paragraphs[0], val, 27, t.font_head, t.ink, bold=True)
        _set(_box(slide, x, top + 0.78, card_w, 1.55).paragraphs[0], desc, 11, t.font_body, GRAY_500)


def _insight_slide_kpi(prs, sec, chart_path: Optional[Path], t):
    """Alternate layout: title on top, a wide chart, then key numbers each with a short
    insight across the bottom — so the deck isn't every-slide chart-left/callout-right."""
    src = f"Source: {sec.citations[0] if sec.citations else 'database'}, verified against the data."
    s = _content_slide(prs, sec.kicker, sec.action_title, src, t)
    if chart_path:
        s.shapes.add_picture(str(chart_path), Inches(1.0), Inches(1.7), width=Inches(11.3))
    _kpi_band(s, _kpis_from_chart(sec.chart), sec.bullets, t)


def _reco_slide(prs, report: Report, t):
    s = _content_slide(prs, "RECOMMENDATIONS", "Priorities for the coming month",
                       "Source: figures verified against the source data.", t)
    tf = _box(s, 0.5, 2.05, CONTENT_W_IN, 4.6)
    for i, rec in enumerate(report.recommendations, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        _set(p, f"{i}.   {rec}", 16, t.font_body, GRAY_900)
        p.space_after = Pt(16)


# --------------------------------------------------------------------------- #
# PDF export
# --------------------------------------------------------------------------- #
def _find_soffice() -> Optional[str]:
    found = shutil.which("soffice") or shutil.which("soffice.com")
    if found:
        return found
    for p in (r"C:\Program Files\LibreOffice\program\soffice.exe",
              r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"):
        if Path(p).exists():
            return p
    return None


def export_pdf(pptx_path: Path, out_dir: Path) -> Optional[Path]:
    soffice = _find_soffice()
    if not soffice:
        return None
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir",
                    str(out_dir), str(pptx_path)], check=True, timeout=180,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = Path(out_dir) / (Path(pptx_path).stem + ".pdf")
    return pdf if pdf.exists() else None


def _save_robust(prs, path: Path) -> Path:
    """Save the deck, falling back to report-1.pptx, -2, ... if the file is open/locked."""
    candidates = [path] + [path.with_name(f"{path.stem}-{i}{path.suffix}") for i in range(1, 50)]
    for cand in candidates:
        try:
            prs.save(str(cand))
            return cand
        except PermissionError:
            continue
    raise PermissionError(f"Could not write {path.name} (is it open in PowerPoint?). "
                          "Close it and try again.")


# --------------------------------------------------------------------------- #
# Orchestration
# --------------------------------------------------------------------------- #
def render_report(report: Report, out_dir=None, charts_dir=None, brand=None) -> dict:
    s = get_settings()
    out_dir = Path(out_dir or s.out_dir)
    charts_dir = Path(charts_dir or s.charts_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    charts_dir.mkdir(parents=True, exist_ok=True)

    # brand (optional) = {accent, ink, font_head, font_body} extracted from an uploaded template.
    t = _theme(brand)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    cover_dark = (brand or {}).get("cover", "dark") != "light"
    _cover_slide(prs, report, t, dark=cover_dark)
    _exec_slide(prs, report, t)
    # Mostly the classic chart-left/takeaway-right layout, with ONE KPI-band slide per
    # deck (the second insight) so there's variety without it being repetitive.
    insight_n = 0
    for i, sec in enumerate(report.sections):
        use_kpi = insight_n == 1 and sec.chart is not None
        cpath = None
        if sec.chart:
            cpath = charts_dir / f"chart_{i}.png"
            render_chart(sec.chart, cpath, accent=t.accent, wide=use_kpi)
        if use_kpi:
            _insight_slide_kpi(prs, sec, cpath, t)
        else:
            _insight_slide(prs, sec, cpath, t)
        insight_n += 1
    _reco_slide(prs, report, t)
    # No standalone "Sources & methodology" slide — a single-line appendix reads as an
    # unfinished slide. Provenance lives in each slide's source footer and in report.json.

    pptx_path = _save_robust(prs, out_dir / "report.pptx")
    pdf_path = export_pdf(pptx_path, out_dir)
    return {"pptx": pptx_path, "pdf": pdf_path, "charts_dir": charts_dir}


def load_report(path) -> Report:
    return Report.model_validate_json(Path(path).read_text(encoding="utf-8"))


if __name__ == "__main__":
    from src.report import build_report
    from src.schemas import ReportRequest

    # Set REUSE_REPORT=1 to re-render the cached report.json without calling the LLM
    # (fast/cheap iteration on the visual design).
    cached = Path(get_settings().out_dir) / "report.json"
    if os.getenv("REUSE_REPORT") == "1" and cached.exists():
        print("(reusing cached report.json — no LLM call)")
        rep = load_report(cached)
    else:
        rep = build_report(ReportRequest(intent="Monthly sales review", table="sales"))
    paths = render_report(rep)
    print("PPTX  :", paths["pptx"])
    print("PDF   :", paths["pdf"] or "(skipped - LibreOffice not installed)")
    print("Charts:", paths["charts_dir"])
